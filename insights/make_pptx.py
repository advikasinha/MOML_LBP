"""
make_pptx.py
============
Generates the MIC-300 Mid-Term presentation PPTX for:
  "Multi-Objective Clustering with Feature Sparsity for Bearing Fault Diagnosis"

Slide structure (strictly follows MIC-300 guidelines):
  Slide 1   : Title — project, course, students, supervisor
  Slide 2   : Problem Definition
  Slide 3   : Research Objectives
  Slide 4   : Methodology
  Slide 5   : Timeline (Gantt)
  Slide 6   : Work Done — Dataset Pipeline & Implementation
  Slide 7   : Work Done — Results & Key Findings
  Slide 8   : Demo (optional) — Interactive Pareto Front

Run:  python make_pptx.py
Output: MIC300_MidTerm_MOC_FS.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── FILL IN YOUR DETAILS ─────────────────────────────────────────────────────
STUDENTS     = ["Advik Goyal", "[Team Member 2]", "[Team Member 3]"]
SUPERVISOR   = "Prof. [Supervisor Name]"
COURSE       = "MIC-300: Lab Based Project  |  Spring 2025-26"
DEPT         = "Department of Mechanical and Industrial Engineering, IIT Roorkee"
DATE         = "March 2026"
# ─────────────────────────────────────────────────────────────────────────────

# Colour palette
C_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
C_ACCENT = RGBColor(0x1B, 0x6C, 0xA8)   # IIT blue
C_GOLD   = RGBColor(0xE8, 0xA8, 0x20)   # warm gold (highlights)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # off-white bg
C_MID    = RGBColor(0xCC, 0xDD, 0xEE)   # light blue
C_RED    = RGBColor(0xC0, 0x39, 0x2B)
C_GREEN  = RGBColor(0x1A, 0x7A, 0x3C)
C_ORANGE = RGBColor(0xD4, 0x6A, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color and line_width:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
             italic=False, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=C_DARK,
             align=PP_ALIGN.LEFT, space_before=6, level=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def section_header(slide, title, subtitle=None, y=0.18, accent=C_ACCENT):
    """Blue left-bar section header."""
    add_rect(slide, 0, y, 0.07, 0.55 if not subtitle else 0.72, fill_color=accent)
    add_text(slide, title, 0.18, y, 12.8, 0.45,
             font_size=26, bold=True, color=C_DARK)
    if subtitle:
        add_text(slide, subtitle, 0.18, y + 0.40, 12.8, 0.35,
                 font_size=14, color=C_ACCENT, italic=True)


def divider(slide, y, color=C_MID):
    add_rect(slide, 0.18, y, 12.95, 0.03, fill_color=color)


def bullet_box(slide, items, left, top, width, height,
               font_size=14, color=C_DARK, bullet="▸", header=None, header_color=None):
    """Box with optional header and bullet items."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    if header:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = header
        run.font.size = Pt(font_size + 1)
        run.font.bold = True
        run.font.color.rgb = header_color or C_ACCENT
        run.font.name = "Calibri"
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return box


def card(slide, left, top, width, height, fill=C_LIGHT, border=C_MID, border_pt=1.2):
    r = add_rect(slide, left, top, width, height, fill_color=fill,
                 line_color=border, line_width=border_pt)
    return r


def metric_card(slide, left, top, w, h, label, value, unit="",
                fill=C_ACCENT, text_color=C_WHITE, val_color=C_GOLD):
    card(slide, left, top, w, h, fill=fill, border=fill)
    add_text(slide, value + unit, left + 0.05, top + 0.05, w - 0.1, h * 0.55,
             font_size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, left + 0.05, top + h * 0.55, w - 0.1, h * 0.42,
             font_size=11, bold=False, color=text_color, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════

def slide1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full dark background
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=C_DARK)

    # Gold top bar
    add_rect(slide, 0, 0, 13.33, 0.12, fill_color=C_GOLD)

    # Blue accent left panel
    add_rect(slide, 0, 0.12, 3.5, 7.38, fill_color=C_ACCENT)

    # Gear / bearing icon area (text representation)
    add_text(slide, "⚙", 0.3, 1.6, 3.0, 2.5,
             font_size=96, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text(slide, "Bearing Fault\nDiagnosis", 0.2, 4.1, 3.1, 1.2,
             font_size=13, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # Course info (top right)
    add_text(slide, COURSE, 3.7, 0.22, 9.3, 0.45,
             font_size=12, color=C_MID, align=PP_ALIGN.LEFT)
    add_text(slide, DEPT, 3.7, 0.58, 9.3, 0.38,
             font_size=11, color=C_MID, italic=True, align=PP_ALIGN.LEFT)

    # Main title
    add_text(slide,
             "Multi-Objective Clustering\nwith Feature Sparsity",
             3.7, 1.2, 9.2, 2.0,
             font_size=38, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    add_text(slide, "for Bearing Fault Diagnosis", 3.7, 3.05, 9.2, 0.65,
             font_size=28, bold=False, color=C_GOLD, align=PP_ALIGN.LEFT)

    # Divider
    add_rect(slide, 3.7, 3.72, 9.2, 0.04, fill_color=C_GOLD)

    # Students & supervisor
    student_str = "  |  ".join(STUDENTS)
    add_text(slide, "Students:", 3.7, 3.85, 1.5, 0.35,
             font_size=12, bold=True, color=C_MID)
    add_text(slide, student_str, 5.0, 3.85, 8.1, 0.35,
             font_size=12, color=C_WHITE)

    add_text(slide, "Supervisor:", 3.7, 4.22, 1.6, 0.35,
             font_size=12, bold=True, color=C_MID)
    add_text(slide, SUPERVISOR, 5.0, 4.22, 8.1, 0.35,
             font_size=12, color=C_WHITE)

    add_text(slide, DATE, 3.7, 4.6, 9.2, 0.35,
             font_size=11, color=C_MID, italic=True)

    # Keywords
    kws = ["NSGA-II", "Unsupervised Learning", "CWRU Dataset", "Pareto Optimisation", "Feature Selection"]
    x = 3.7
    for kw in kws:
        w = len(kw) * 0.115 + 0.25
        r = add_rect(slide, x, 5.15, w, 0.38,
                     fill_color=RGBColor(0x1B, 0x4F, 0x72), line_color=C_ACCENT, line_width=1)
        add_text(slide, kw, x + 0.06, 5.18, w - 0.08, 0.32,
                 font_size=11, color=C_WHITE, align=PP_ALIGN.CENTER)
        x += w + 0.15

    # IIT Roorkee bottom left
    add_text(slide, "IIT Roorkee", 0.1, 6.85, 3.3, 0.5,
             font_size=13, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM DEFINITION
# ═════════════════════════════════════════════════════════════════════════════

def slide2_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=C_LIGHT)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=C_GOLD)

    section_header(slide, "Problem Definition", "Why do bearing faults matter — and why is this hard?")
    divider(slide, 0.92)

    # --- Stat cards row ---
    stats = [
        ("40–50%", "of rotating machine\nfailures from bearings"),
        ("4 Fault\nTypes",  "Normal · IR · Ball · OR\n@ 4 severity levels"),
        ("4 Load\nConds.",  "0–3 HP  |  1730–1797 RPM\nmasks fault signatures"),
        ("7 Features\nExtracted", "RMS, Kurtosis, Skewness,\nCrest, P2P, Std, Spec.Centroid"),
    ]
    colors = [C_DARK, C_ACCENT, C_RED, C_GREEN]
    for i, (val, lab) in enumerate(stats):
        x = 0.18 + i * 3.27
        metric_card(slide, x, 1.08, 3.0, 1.28, lab, val, fill=colors[i])

    # --- Challenge boxes ---
    challenges = [
        ("Challenge 1: Load-Condition Artefacts",
         C_RED,
         ["Different motor loads (0–3 HP) shift vibration amplitude",
          "Global normalisation fuses load variation with fault signatures",
          "Fix: per-load-condition StandardScaler before merging"]),
        ("Challenge 2: Class Imbalance",
         C_ORANGE,
         ["OR fault has 3 mounting positions (@3, @6, @12) → 3× more data",
          "Biases the clustering objective toward the over-represented class",
          "Fix: Random undersampling to median class count"]),
        ("Challenge 3: No Ground-Truth for Clustering",
         C_ACCENT,
         ["Bearing fault diagnosis is inherently unsupervised at deployment",
          "Labels only available for validation — not used during optimisation",
          "Evaluation: Adjusted Rand Index (ARI) + Silhouette Score"]),
        ("Challenge 4: Competing Objectives",
         C_GREEN,
         ["Tight clusters ≠ well-separated classes (compactness vs. separation)",
          "More features ≠ better clusters (curse of dimensionality)",
          "Solution: 3-objective Pareto optimisation (f1, f2, f3)"]),
    ]

    for i, (title, color, pts) in enumerate(challenges):
        row, col = divmod(i, 2)
        x = 0.18 + col * 6.55
        y = 2.55 + row * 2.28
        card(slide, x, y, 6.3, 2.15, fill=C_WHITE, border=color, border_pt=2.0)
        add_rect(slide, x, y, 6.3, 0.4, fill_color=color)
        add_text(slide, title, x + 0.12, y + 0.05, 6.0, 0.35,
                 font_size=12, bold=True, color=C_WHITE)
        for j, pt in enumerate(pts):
            add_text(slide, f"▸  {pt}", x + 0.12, y + 0.48 + j * 0.52, 6.0, 0.48,
                     font_size=11.5, color=C_DARK)

    # Slide number
    add_text(slide, "2", 12.9, 7.1, 0.4, 0.35, font_size=10, color=C_ACCENT, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — RESEARCH OBJECTIVES
# ═════════════════════════════════════════════════════════════════════════════

def slide3_objectives(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=C_LIGHT)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=C_GOLD)

    section_header(slide, "Research Objectives",
                   "A rigorous multi-objective framework for unsupervised fault discovery")
    divider(slide, 0.92)

    objectives = [
        ("RO 1", "Dataset Pipeline",
         C_ACCENT,
         "Build a robust preprocessing pipeline:\n"
         "per-load-condition normalisation  ·  class balancing  ·  16-class severity labelling\n"
         "→  7,335 balanced samples from 12k Drive-End CWRU dataset"),
        ("RO 2", "Chromosome Encoding",
         C_GREEN,
         "Design a unified variable-length chromosome:\n"
         "[ K ]  [ K×D cluster centres ]  [ D-bit feature mask ]\n"
         "→  78 variables representing K ∈ {2…10}, D = 7 features, simultaneous K + feature selection"),
        ("RO 3", "Three-Objective NSGA-II",
         C_RED,
         "Simultaneously minimise:\n"
         "f1 Compactness (TWCSS/N/d)  ·  f2 Connectedness (kNN cross-cluster fraction)  ·  f3 Simplicity (L0 feature count)\n"
         "→  200-population × 200-generation NSGA-II with hybrid initialisation"),
        ("RO 4", "Pareto Analysis & Validation",
         C_ORANGE,
         "Extract Pareto front  →  knee-point selection via normalised Chebyshev distance\n"
         "Validate against ground-truth: ARI (4-class fault type)  +  ARI (16-class type × severity)\n"
         "→  Silhouette Score as internal cohesion metric"),
        ("RO 5", "Physical Interpretation",
         C_DARK,
         "Map discovered minimal feature subsets to bearing fault physics\n"
         "Identify which vibration statistics (RMS, Kurtosis, Spectral Centroid…) catch which fault modes\n"
         "→  BPFI / BPFO / BSF marker validation via FFT analysis"),
    ]

    for i, (tag, name, color, desc) in enumerate(objectives):
        y = 1.08 + i * 1.22
        # Number bubble
        add_rect(slide, 0.18, y + 0.05, 0.55, 0.55, fill_color=color)
        add_text(slide, tag, 0.18, y + 0.05, 0.55, 0.55,
                 font_size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Card
        card(slide, 0.8, y, 12.38, 1.08, fill=C_WHITE, border=color, border_pt=1.5)
        add_text(slide, name, 0.95, y + 0.06, 11.9, 0.38,
                 font_size=13, bold=True, color=color)
        add_text(slide, desc, 0.95, y + 0.38, 12.0, 0.65,
                 font_size=11, color=C_DARK)

    add_text(slide, "3", 12.9, 7.1, 0.4, 0.35, font_size=10, color=C_ACCENT, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — METHODOLOGY
# ═════════════════════════════════════════════════════════════════════════════

def slide4_methodology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=C_LIGHT)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=C_GOLD)

    section_header(slide, "Methodology", "Three-stage pipeline: Data → Optimisation → Interpretation")
    divider(slide, 0.92)

    # ── Stage 1: Data Pipeline ────────────────────────────────────────────
    add_rect(slide, 0.18, 1.05, 4.0, 0.42, fill_color=C_ACCENT)
    add_text(slide, "① Data Pipeline", 0.28, 1.08, 3.8, 0.38,
             font_size=13, bold=True, color=C_WHITE)

    pipe_steps = [
        ("CWRU 12kHz .mat files",           "12k Drive-End, 4 fault types, 4 sizes, 4 loads"),
        ("Segment → Windows (n=1024)",       "Extract fixed-length segments, compute 7 features/window"),
        ("Per-load StandardScaler",          "Separate normaliser per load condition (FIX-1)"),
        ("Random Undersampling",             "Balance to median class count (FIX-2)"),
        ("16-class composite label",         "fault_type × fault_size — for ARI validation (FIX-3)"),
    ]
    for j, (step, detail) in enumerate(pipe_steps):
        y = 1.52 + j * 0.46
        add_rect(slide, 0.18, y, 0.08, 0.35, fill_color=C_GOLD)
        add_text(slide, step, 0.32, y + 0.01, 1.85, 0.35, font_size=10.5, bold=True, color=C_DARK)
        add_text(slide, detail, 2.18, y + 0.01, 2.0, 0.35, font_size=10, color=RGBColor(0x44,0x44,0x44))

    # ── Arrow between stages ──────────────────────────────────────────────
    add_rect(slide, 4.22, 1.3, 0.5, 3.7, fill_color=C_MID)
    add_text(slide, "▶", 4.22, 2.9, 0.5, 0.5, font_size=18, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # ── Stage 2: NSGA-II ─────────────────────────────────────────────────
    add_rect(slide, 4.85, 1.05, 4.85, 0.42, fill_color=C_RED)
    add_text(slide, "② NSGA-II Optimisation", 4.95, 1.08, 4.65, 0.38,
             font_size=13, bold=True, color=C_WHITE)

    nsga_items = [
        ("Chromosome", "[ K | C₁…C_K×D | mask₁…mask_D ]   (78 vars)"),
        ("Population", "200 individuals  |  Hybrid init: KMeans + random"),
        ("Objectives", "f1 TWCSS/N/d  ·  f2 kNN-penalty  ·  f3 L0-norm"),
        ("Operators",  "SBX + 1-pt crossover  |  Polynomial + bit-flip mutation"),
        ("Budget",     "200 generations  →  Pareto front (200 solutions)"),
    ]
    for j, (key, val) in enumerate(nsga_items):
        y = 1.52 + j * 0.46
        add_text(slide, key + ":", 4.85, y, 1.2, 0.38,
                 font_size=10.5, bold=True, color=C_RED)
        add_text(slide, val, 6.05, y, 3.55, 0.38,
                 font_size=10.5, color=C_DARK)

    # ── Arrow ─────────────────────────────────────────────────────────────
    add_rect(slide, 9.74, 1.3, 0.5, 3.7, fill_color=C_MID)
    add_text(slide, "▶", 9.74, 2.9, 0.5, 0.5, font_size=18, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # ── Stage 3: Analysis ─────────────────────────────────────────────────
    add_rect(slide, 10.28, 1.05, 2.9, 0.42, fill_color=C_GREEN)
    add_text(slide, "③ Analysis", 10.38, 1.08, 2.7, 0.38,
             font_size=13, bold=True, color=C_WHITE)

    analysis_items = [
        ("Knee point",    "Min Chebyshev dist.\nto ideal in norm. space"),
        ("Validation",    "ARI (4-cls + 16-cls)\nSilhouette Score"),
        ("Interpretation","Feature mask → physics\nFFT + fault freq. markers"),
    ]
    for j, (key, val) in enumerate(analysis_items):
        y = 1.52 + j * 0.62
        card(slide, 10.28, y, 2.9, 0.55, fill=C_WHITE, border=C_GREEN, border_pt=1.2)
        add_text(slide, key, 10.38, y + 0.04, 2.7, 0.25,
                 font_size=11, bold=True, color=C_GREEN)
        add_text(slide, val, 10.38, y + 0.26, 2.7, 0.28,
                 font_size=10, color=C_DARK)

    # ── Objective equations ───────────────────────────────────────────────
    add_rect(slide, 0.18, 4.82, 12.95, 0.04, fill_color=C_MID)
    add_rect(slide, 0.18, 4.90, 12.95, 2.48, fill_color=C_WHITE,
             line_color=C_MID, line_width=1)
    add_text(slide, "Three Objectives Simultaneously Optimised by NSGA-II",
             0.35, 4.94, 12.5, 0.38, font_size=12, bold=True, color=C_DARK)

    eqs = [
        ("f₁  Compactness", "min  TWCSS / (N · d_active)",
         "Total within-cluster sum of squares — measures how tightly grouped the data is",
         C_ACCENT),
        ("f₂  Connectedness", "min  Σ penalty(xᵢ, kNN(xᵢ))",
         "Fraction of k nearest-neighbours that lie in a different cluster",
         C_RED),
        ("f₃  Simplicity", "min  ‖mask‖₀  ∈ {1…7}",
         "L0-norm of feature mask — forces the algorithm to use as few features as possible",
         C_GREEN),
    ]
    for i, (name, eq, desc, color) in enumerate(eqs):
        x = 0.35 + i * 4.3
        add_text(slide, name, x, 5.35, 4.0, 0.32, font_size=11, bold=True, color=color)
        add_text(slide, eq,   x, 5.65, 4.0, 0.35, font_size=13, bold=True, color=C_DARK,
                 font_name="Courier New")
        add_text(slide, desc, x, 5.98, 4.1, 0.55, font_size=10, color=RGBColor(0x44,0x44,0x44))

    add_text(slide, "4", 12.9, 7.1, 0.4, 0.35, font_size=10, color=C_ACCENT, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TIMELINE (GANTT)
# ═════════════════════════════════════════════════════════════════════════════

def slide5_timeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=C_LIGHT)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=C_GOLD)

    section_header(slide, "Project Timeline", "MIC-300  |  Semester: November 2025 – May 2026")
    divider(slide, 0.92)

    # Months header
    months = ["Nov '25", "Dec '25", "Jan '26", "Feb '26", "Mar '26", "Apr '26", "May '26"]
    col_start = 3.2
    col_w     = 1.38
    row_h     = 0.62

    for i, m in enumerate(months):
        x = col_start + i * col_w
        bg = C_ACCENT if i in [3, 4] else C_DARK  # highlight current months
        add_rect(slide, x, 1.05, col_w - 0.05, 0.42, fill_color=bg)
        add_text(slide, m, x, 1.08, col_w - 0.05, 0.38,
                 font_size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Tasks
    tasks = [
        # (label, start_month_idx, end_month_idx, color, status)
        ("Literature Review & Dataset Acquisition",      0, 1,  C_DARK,   "Done"),
        ("CWRU Preprocessing Pipeline (FIX 1-3)",        1, 2,  C_GREEN,  "Done"),
        ("Chromosome Encoding & NSGA-II Framework",      1, 3,  C_ACCENT, "Done"),
        ("v1 Pipeline Run & Pareto Analysis",            2, 3,  C_ACCENT, "Done"),
        ("v2 Bug Fixes & Full Dataset Re-run",           3, 4,  C_GREEN,  "Done ✓"),
        ("Mechanical Interpretation & Visualisations",   3, 4,  C_ORANGE, "Done ✓"),
        ("Comparative Analysis & Insights Documentation",4, 5,  C_RED,    "In Progress"),
        ("End-term Write-up & Final Report",             5, 6,  C_DARK,   "Upcoming"),
    ]

    for j, (label, s, e, color, status) in enumerate(tasks):
        y = 1.52 + j * row_h
        # Row background
        bg_row = C_WHITE if j % 2 == 0 else RGBColor(0xE8, 0xF2, 0xFB)
        add_rect(slide, 0.18, y, 13.0, row_h - 0.06, fill_color=bg_row)
        # Task label
        add_text(slide, label, 0.25, y + 0.05, 2.9, row_h - 0.14,
                 font_size=10.5, color=C_DARK)
        # Gantt bar
        bar_x = col_start + s * col_w
        bar_w = (e - s + 1) * col_w - 0.12
        add_rect(slide, bar_x, y + 0.1, bar_w, row_h - 0.26, fill_color=color)
        # Status badge
        status_color = C_GREEN if "Done" in status else (C_ORANGE if "Progress" in status else C_MID)
        add_text(slide, status, bar_x + bar_w * 0.5 - 0.4, y + 0.12, 1.2, 0.35,
                 font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Legend
    for label, color in [("Completed", C_GREEN), ("In Progress", C_ORANGE),
                          ("Upcoming", C_DARK), ("Current Period", C_ACCENT)]:
        x_leg = 0.25 + [("Completed", C_GREEN), ("In Progress", C_ORANGE),
                         ("Upcoming", C_DARK), ("Current Period", C_ACCENT)].index((label, color)) * 3.0
        add_rect(slide, x_leg, 7.1, 0.25, 0.22, fill_color=color)
        add_text(slide, label, x_leg + 0.3, 7.08, 2.6, 0.28, font_size=10, color=C_DARK)

    add_text(slide, "5", 12.9, 7.1, 0.4, 0.35, font_size=10, color=C_ACCENT, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — WORK DONE SO FAR (Dataset & Implementation)
# ═════════════════════════════════════════════════════════════════════════════

def slide6_work_done(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=C_LIGHT)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=C_GOLD)

    section_header(slide, "Work Done So Far — Part 1",
                   "Dataset preparation, bug fixes and algorithm implementation")
    divider(slide, 0.92)

    # ── Dataset stats ─────────────────────────────────────────────────────
    add_rect(slide, 0.18, 1.05, 12.95, 0.38, fill_color=C_ACCENT)
    add_text(slide, "A  CWRU Dataset v2 — Three Critical Fixes Applied",
             0.28, 1.07, 12.7, 0.34, font_size=13, bold=True, color=C_WHITE)

    stats = [
        ("7,335",  "Balanced\nsamples"),
        ("12",     "Classes\n(16-class)"),
        ("7",      "Features\nextracted"),
        ("4",      "Fault\ntypes"),
        ("28",     "OR files\nloaded"),
        ("4",      "Load\nconditions"),
    ]
    for i, (val, lab) in enumerate(stats):
        x = 0.18 + i * 2.15
        metric_card(slide, x, 1.48, 2.0, 1.0, lab, val,
                    fill=C_DARK if i % 2 == 0 else C_ACCENT)

    # Fix cards
    fixes = [
        ("FIX 1: Per-Load Normalisation", C_RED,
         ["Old: single global StandardScaler mixed load variation with fault signal",
          "New: separate scaler per load condition (0, 1, 2, 3 HP) before merging",
          "Impact: primary driver of ARI improvement — load artefacts no longer dominate"]),
        ("FIX 2: Class Balancing", C_ORANGE,
         ["OR fault had 3 mounting positions → 3× more segments than other classes",
          "Random undersampling to median class count across all 4 fault types",
          "Impact: prevents the clustering from being biased towards OR"]),
        ("FIX 3: OR Recursive Glob + 16-class Labels", C_GREEN,
         ["Bug: non-recursive glob missed OR/007/@3/, @6/, @12/ subdirectories",
          "Fix: glob(..., recursive=True) now loads all 28 OR files (was 4)",
          "Severity label: fault_type × fault_size → 12 distinct classes for ARI-16"]),
    ]
    for i, (title, color, pts) in enumerate(fixes):
        x = 0.18 + i * 4.35
        y = 2.58
        card(slide, x, y, 4.12, 2.35, fill=C_WHITE, border=color, border_pt=2.0)
        add_rect(slide, x, y, 4.12, 0.38, fill_color=color)
        add_text(slide, title, x + 0.1, y + 0.05, 3.9, 0.32,
                 font_size=11.5, bold=True, color=C_WHITE)
        for j, pt in enumerate(pts):
            add_text(slide, f"▸  {pt}", x + 0.1, y + 0.44 + j * 0.62, 3.9, 0.57,
                     font_size=10.5, color=C_DARK)

    # ── Implementation highlights ─────────────────────────────────────────
    add_rect(slide, 0.18, 5.05, 12.95, 0.38, fill_color=C_DARK)
    add_text(slide, "B  NSGA-II Implementation Highlights",
             0.28, 5.07, 12.7, 0.34, font_size=13, bold=True, color=C_WHITE)

    impl = [
        ("Hybrid Init",      "1/3 KMeans-seeded  +  1/3 random data-point centres  +  1/3 pure random"),
        ("Custom Crossover", "SBX on continuous centre genes  +  1-point crossover on 7-bit feature mask"),
        ("Custom Mutation",  "Polynomial mutation + bit-flip (p=0.15) + structural K±1 (p=0.1)"),
        ("History Callbacks","Snapshots at gens {1,5,10,20,35,50,75,100,150,200} for convergence study"),
    ]
    for j, (key, val) in enumerate(impl):
        y = 5.48 + j * 0.44
        add_rect(slide, 0.18, y, 2.1, 0.38, fill_color=C_ACCENT)
        add_text(slide, key, 0.22, y + 0.05, 2.0, 0.3, font_size=10.5, bold=True, color=C_WHITE)
        add_text(slide, val, 2.32, y + 0.05, 10.7, 0.38, font_size=10.5, color=C_DARK)

    add_text(slide, "6", 12.9, 7.1, 0.4, 0.35, font_size=10, color=C_ACCENT, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WORK DONE SO FAR (Results)
# ═════════════════════════════════════════════════════════════════════════════

def slide7_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=C_LIGHT)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=C_GOLD)

    section_header(slide, "Work Done So Far — Part 2",
                   "NSGA-II results, Pareto front analysis and physical interpretation")
    divider(slide, 0.92)

    # ── Key metric cards ──────────────────────────────────────────────────
    metrics = [
        ("0.36",  "ARI (4-class)\nfault type", C_ACCENT),
        ("0.50",  "ARI (16-class)\ntype×severity", C_GREEN),
        ("0.82",  "Silhouette\nScore", C_RED),
        ("200",   "Pareto\nSolutions", C_DARK),
        ("1–2",   "f3 active\nfeatures", C_ORANGE),
        ("RMS",   "Dominant\nfeature", C_DARK),
    ]
    for i, (val, lab, color) in enumerate(metrics):
        x = 0.18 + i * 2.17
        metric_card(slide, x, 1.05, 2.0, 1.1, lab, val, fill=color)

    # ── Left: Pareto analysis table ───────────────────────────────────────
    add_rect(slide, 0.18, 2.28, 6.2, 0.40, fill_color=C_ACCENT)
    add_text(slide, "Pareto Front Analysis", 0.28, 2.30, 6.0, 0.36,
             font_size=13, bold=True, color=C_WHITE)

    rows = [
        ("Solution type",    "K",  "Features",         "ARI (4cls)"),
        ("Knee point",       "10", "RMS",              "0.3595"),
        ("Best ARI",         "10", "RMS",              "0.3595"),
        ("Most compact",     "10", "RMS",              "0.2904"),
        ("f3=2 example",     " 5", "RMS + Kurtosis",   "0.3102"),
    ]
    row_colors = [C_DARK, C_WHITE, C_WHITE, C_WHITE, C_WHITE]
    text_colors= [C_WHITE, C_DARK, C_DARK, C_DARK, C_DARK]
    for j, (r1, r2, r3, r4) in enumerate(rows):
        y = 2.72 + j * 0.52
        add_rect(slide, 0.18, y, 6.2, 0.5, fill_color=row_colors[j],
                 line_color=C_MID, line_width=0.8)
        for k, txt in enumerate([r1, r2, r3, r4]):
            add_text(slide, txt, 0.22 + k * 1.52, y + 0.1, 1.48, 0.38,
                     font_size=11, bold=(j == 0), color=text_colors[j])

    # ── Right: Physical interpretation ───────────────────────────────────
    add_rect(slide, 6.55, 2.28, 6.6, 0.40, fill_color=C_GREEN)
    add_text(slide, "Physical Interpretation", 6.65, 2.30, 6.4, 0.36,
             font_size=13, bold=True, color=C_WHITE)

    phys = [
        ("▸  RMS dominates",
         "After per-load normalisation, overall energy level (RMS) is the\n"
         "single most discriminative feature — confirming that fault-induced\n"
         "energy changes are the primary clustering signal."),
        ("▸  f3 Pareto range: 1–2 features",
         "Adding features beyond RMS does not improve cluster quality.\n"
         "Curse of dimensionality: 7 features all similarly discriminative\n"
         "→ sparse is better in this 7-dimensional space."),
        ("▸  16-class ARI > 4-class ARI (0.50 vs 0.36)",
         "Clusters align more closely with fault_type × severity than\n"
         "fault_type alone — the algorithm discovers severity sub-structure\n"
         "even without severity labels during training."),
    ]
    for j, (hdr, body) in enumerate(phys):
        y = 2.72 + j * 1.42
        card(slide, 6.55, y, 6.6, 1.35, fill=C_WHITE,
             border=C_GREEN, border_pt=1.5)
        add_text(slide, hdr, 6.65, y + 0.06, 6.38, 0.35,
                 font_size=11.5, bold=True, color=C_GREEN)
        add_text(slide, body, 6.65, y + 0.38, 6.38, 0.92,
                 font_size=10.5, color=C_DARK)

    # ── Bottom: v1 vs v2 comparison ───────────────────────────────────────
    add_rect(slide, 0.18, 6.48, 12.95, 0.04, fill_color=C_MID)
    add_rect(slide, 0.18, 6.55, 12.95, 0.82, fill_color=C_WHITE,
             line_color=C_MID, line_width=1)
    add_text(slide, "v1 → v2 Comparison:", 0.3, 6.60, 2.3, 0.35,
             font_size=11, bold=True, color=C_DARK)
    add_text(slide,
             "Dataset: 5,678 → 7,335 samples   |   OR files loaded: 4 → 28   |   "
             "Silhouette: 0.838 → 0.815   |   "
             "Note: lower ARI in v2 reflects harder, more realistic clustering (full OR class variation)",
             2.5, 6.60, 10.7, 0.68, font_size=10.5, color=C_DARK)

    add_text(slide, "7", 12.9, 7.1, 0.4, 0.35, font_size=10, color=C_ACCENT, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DEMO (optional)
# ═════════════════════════════════════════════════════════════════════════════

def slide8_demo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=C_DARK)
    add_rect(slide, 0, 0, 13.33, 0.08, fill_color=C_GOLD)

    # Header
    add_rect(slide, 0, 0.08, 13.33, 0.78, fill_color=C_ACCENT)
    add_text(slide, "Live Demonstration  —  Interactive 3-D Pareto Front",
             0.25, 0.12, 12.8, 0.66, font_size=22, bold=True, color=C_WHITE)

    # Main info box
    card(slide, 0.5, 1.0, 12.33, 5.6, fill=RGBColor(0x1A, 0x2A, 0x3A),
         border=C_ACCENT, border_pt=2)

    add_text(slide,
             "pareto_3d_interactive.html",
             0.7, 1.15, 11.9, 0.55,
             font_size=18, bold=True, color=C_GOLD, font_name="Courier New")

    points = [
        ("Drag to rotate",        "Fully rotatable 3-D scatter of all 200 Pareto solutions"),
        ("Hover for details",     "Each point shows: K, active features, ARI, Silhouette, f1/f2/f3 values"),
        ("Red ◆ = Knee point",   "Minimum Chebyshev distance to ideal — best balanced solution"),
        ("Gold ■ = Best ARI",    "Solution with highest Adjusted Rand Index vs. fault labels"),
        ("Colour = ARI",          "Green → higher ARI (better fault alignment), Red → lower ARI"),
        ("Axes",                  "f1 Compactness  ×  f2 Connectedness  ×  f3 Simplicity (# features)"),
    ]
    for j, (key, val) in enumerate(points):
        y = 1.82 + j * 0.7
        add_rect(slide, 0.65, y, 2.3, 0.52, fill_color=C_ACCENT)
        add_text(slide, key, 0.7, y + 0.08, 2.2, 0.38,
                 font_size=11, bold=True, color=C_WHITE)
        add_text(slide, val, 3.02, y + 0.08, 9.5, 0.52,
                 font_size=11.5, color=C_WHITE)

    add_text(slide,
             "Open  moc_results_v2/pareto_3d_interactive.html  in any browser — no server required",
             0.7, 6.15, 12.0, 0.4,
             font_size=11, italic=True, color=C_MID, font_name="Courier New")

    add_text(slide, "8", 12.9, 7.1, 0.4, 0.35, font_size=10, color=C_GOLD, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides…")
    slide1_title(prs)       ; print("  [1/8] Title")
    slide2_problem(prs)     ; print("  [2/8] Problem Definition")
    slide3_objectives(prs)  ; print("  [3/8] Research Objectives")
    slide4_methodology(prs) ; print("  [4/8] Methodology")
    slide5_timeline(prs)    ; print("  [5/8] Timeline")
    slide6_work_done(prs)   ; print("  [6/8] Work Done – Part 1")
    slide7_results(prs)     ; print("  [7/8] Work Done – Part 2 (Results)")
    slide8_demo(prs)        ; print("  [8/8] Demo")

    out = "MIC300_MidTerm_MOC_FS.pptx"
    prs.save(out)
    print(f"\nSaved: {out}")
    print("→  Upload to Google Slides:  slides.google.com  →  File → Import slides")
    print("\nREMEMBER to fill in:")
    print("  STUDENTS   =", STUDENTS)
    print("  SUPERVISOR =", SUPERVISOR)


if __name__ == "__main__":
    main()
